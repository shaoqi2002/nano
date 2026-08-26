from pathlib import Path
from typing import Any
from io import BytesIO
from zipfile import BadZipFile, ZipFile
from dataclasses import dataclass
import colorsys
import hashlib
import json
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


NAVY = RGBColor(0x13, 0x26, 0x3D)
BLUE = RGBColor(0x2E, 0x74, 0xB5)
LIGHT_BLUE = RGBColor(0xE9, 0xF2, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x20, 0x2A, 0x35)
GRAY = RGBColor(0x67, 0x72, 0x7E)
LIGHT_GRAY = RGBColor(0xF3, 0xF5, 0xF7)


@dataclass(frozen=True)
class PresentationTheme:
    name: str
    background: RGBColor
    surface: RGBColor
    primary: RGBColor
    accent: RGBColor
    text: RGBColor
    muted: RGBColor
    subtle: RGBColor
    on_primary: RGBColor = WHITE


THEMES = {
    "business": PresentationTheme(
        "business", RGBColor(0xF8, 0xFA, 0xFC), WHITE, NAVY, BLUE, DARK, GRAY,
        LIGHT_BLUE,
    ),
    "modern": PresentationTheme(
        "modern", RGBColor(0xFB, 0xF7, 0xF2), RGBColor(0xFF, 0xFC, 0xF8),
        RGBColor(0x2D, 0x2A, 0x32), RGBColor(0xE8, 0x68, 0x4A),
        RGBColor(0x2D, 0x2A, 0x32), RGBColor(0x76, 0x6E, 0x68),
        RGBColor(0xF7, 0xE4, 0xDC),
    ),
    "tech": PresentationTheme(
        "tech", RGBColor(0x0D, 0x16, 0x25), RGBColor(0x15, 0x23, 0x38),
        RGBColor(0x0D, 0x16, 0x25), RGBColor(0x38, 0xBD, 0xF8), WHITE,
        RGBColor(0xA6, 0xB3, 0xC5), RGBColor(0x1E, 0x3A, 0x5F),
    ),
}


def resolve_theme(name: str | None) -> PresentationTheme:
    normalized = str(name or "business").strip().lower()
    aliases = {"dark": "tech", "technology": "tech", "warm": "modern", "商务": "business", "科技": "tech", "简约": "modern"}
    normalized = aliases.get(normalized, normalized)
    if normalized not in THEMES:
        raise PresentationArtifactError(
            f"不支持的 PPT 主题：{name}；可选 business、modern、tech"
        )
    return THEMES[normalized]


def _rgb_from_hls(hue: float, lightness: float, saturation: float) -> RGBColor:
    red, green, blue = colorsys.hls_to_rgb(hue % 1.0, lightness, saturation)
    return RGBColor(round(red * 255), round(green * 255), round(blue * 255))


def _parse_hex_color(value: Any, field: str) -> RGBColor | None:
    if value in (None, ""):
        return None
    match = re.fullmatch(r"#?([0-9a-fA-F]{6})", str(value).strip())
    if not match:
        raise PresentationArtifactError(f"{field} 必须是 #RRGGBB 格式")
    raw = match.group(1)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _relative_luminance(color: RGBColor) -> float:
    channels = []
    for value in color:
        channel = value / 255
        channels.append(channel / 12.92 if channel <= 0.04045 else ((channel + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast_ratio(first: RGBColor, second: RGBColor) -> float:
    brighter, darker = sorted(
        (_relative_luminance(first), _relative_luminance(second)), reverse=True
    )
    return (brighter + 0.05) / (darker + 0.05)


def build_adaptive_theme(
    title: str,
    subtitle: str,
    slides: list[dict[str, Any]],
    design_profile: dict[str, Any] | None = None,
) -> PresentationTheme:
    profile = dict(design_profile or {})
    mood = str(profile.get("mood") or profile.get("style") or "").strip().lower()
    structural_content = [
        {
            "title": _first_value(slide, "title", "heading", "name"),
            "kind": _first_value(slide, "type", "layout", "kind"),
            "summary": str(_first_value(slide, "content", "body", "text") or "")[:120],
        }
        for slide in slides
    ]
    serialized = json.dumps(structural_content, ensure_ascii=False, sort_keys=True, default=str)
    digest = hashlib.sha256(f"{title}\n{subtitle}\n{serialized}\n{mood}".encode("utf-8")).digest()
    hue = int.from_bytes(digest[:2], "big") / 65535
    saturation = 0.48 + digest[2] / 255 * 0.18
    accent_offset = 0.08 + digest[3] / 255 * 0.13
    requested_mode = str(profile.get("mode") or profile.get("background_mode") or "").lower()
    dark_mode = requested_mode == "dark" or any(
        token in mood for token in ("dark", "深色", "暗色", "夜间", "cinematic")
    )

    primary = _parse_hex_color(profile.get("primary_color"), "primary_color")
    accent = _parse_hex_color(profile.get("accent_color"), "accent_color")
    background = _parse_hex_color(profile.get("background_color"), "background_color")
    text = _parse_hex_color(profile.get("text_color"), "text_color")
    primary = primary or _rgb_from_hls(hue, 0.24 if not dark_mode else 0.16, saturation)
    accent = accent or _rgb_from_hls(hue + accent_offset, 0.53 if not dark_mode else 0.62, min(0.82, saturation + 0.14))
    if dark_mode:
        background = background or _rgb_from_hls(hue, 0.075, 0.34)
        surface = _rgb_from_hls(hue, 0.12, 0.3)
        text = text or RGBColor(0xF7, 0xFA, 0xFC)
        muted = _rgb_from_hls(hue, 0.72, 0.18)
        subtle = _rgb_from_hls(hue + accent_offset, 0.22, 0.38)
    else:
        background = background or _rgb_from_hls(hue, 0.975, 0.2)
        surface = _rgb_from_hls(hue, 0.995, 0.08)
        text = text or _rgb_from_hls(hue, 0.16, 0.28)
        muted = _rgb_from_hls(hue, 0.43, 0.16)
        subtle = _rgb_from_hls(hue + accent_offset, 0.91, 0.38)
    if _contrast_ratio(text, background) < 4.5:
        if profile.get("text_color"):
            raise PresentationArtifactError("text_color 与 background_color 的对比度不足")
        text = WHITE if _relative_luminance(background) < 0.35 else DARK
    if profile.get("accent_color") and _contrast_ratio(accent, background) < 2.2:
        raise PresentationArtifactError("accent_color 与背景过于接近")
    on_primary = DARK if _relative_luminance(primary) > 0.48 else WHITE
    theme_id = hashlib.sha256(bytes(primary) + bytes(accent) + bytes(background)).hexdigest()[:8]
    return PresentationTheme(
        name=f"adaptive-{theme_id}",
        background=background,
        surface=surface,
        primary=primary,
        accent=accent,
        text=text,
        muted=muted,
        subtle=subtle,
        on_primary=on_primary,
    )


def select_theme_for_content(
    title: str,
    subtitle: str,
    slides: list[dict[str, Any]],
    design_profile: dict[str, Any] | None = None,
) -> PresentationTheme:
    return build_adaptive_theme(title, subtitle, slides, design_profile)


class PresentationArtifactError(RuntimeError):
    pass


def presentation_text(content: bytes) -> str:
    try:
        with ZipFile(BytesIO(content)) as archive:
            entries = archive.infolist()
            if len(entries) > 10_000 or sum(item.file_size for item in entries) > 80 * 1024 * 1024:
                raise PresentationArtifactError("PPTX 解压后内容过大")
            if "ppt/presentation.xml" not in {item.filename for item in entries}:
                raise PresentationArtifactError("PPTX 缺少演示文稿内容")
        presentation = Presentation(BytesIO(content))
    except PresentationArtifactError:
        raise
    except BadZipFile as error:
        raise PresentationArtifactError("PPTX 文件无法解析") from error
    except Exception as error:
        raise PresentationArtifactError("PPTX 文件无法解析") from error
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    blocks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    blocks.append("\t".join(cell.text.strip() for cell in row.cells))
        if blocks:
            slides.append(f"[幻灯片 {index}]\n" + "\n".join(blocks))
    return "\n\n".join(slides)


def _set_run_font(run, size: float, *, bold: bool = False, color=DARK) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: float = 20,
    bold: bool = False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = str(text)
    _set_run_font(run, size, bold=bold, color=color)
    return shape


def _add_slide_header(slide, title: str, number: int, theme: PresentationTheme) -> None:
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent
    accent.line.fill.background()
    title_size = 36 if len(title) <= 28 else 32
    _add_textbox(slide, 0.72, 0.42, 11.35, 0.72, title, size=title_size, bold=True, color=theme.text)
    _add_textbox(
        slide, 12.15, 0.57, 0.55, 0.35, f"{number:02d}", size=11, color=theme.muted,
        align=PP_ALIGN.RIGHT,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.18), Inches(12.0), Inches(0.025)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.subtle
    line.line.fill.background()


def _add_bullets(frame, items: list[Any], *, numbered: bool = False, font_size: float = 20) -> None:
    frame.clear()
    frame.word_wrap = True
    if not items:
        items = [""]
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        if isinstance(item, dict):
            text = str(item.get("text") or "")
            level = min(max(int(item.get("level", 0)), 0), 4)
        else:
            text = str(item)
            level = 0
        if numbered and level == 0:
            paragraph.text = f"{index + 1}. {text}"
        else:
            paragraph.text = f"{'  ' * level}• {text}"
        paragraph.level = level
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = 1.12
        for run in paragraph.runs:
            _set_run_font(run, max(15, font_size - level * 2), color=DARK)


def _add_content_block(
    slide, block: dict[str, Any], left: float, top: float, width: float,
    theme: PresentationTheme,
) -> float:
    kind = str(block.get("type") or "paragraph")
    if kind == "heading":
        _add_textbox(slide, left, top, width, 0.5, str(block.get("text") or ""), size=24, bold=True, color=theme.accent)
        return 0.62
    if kind == "paragraph":
        text = str(block.get("text") or "")
        height = min(1.55, max(0.48, 0.34 + len(text) / 115))
        _add_textbox(slide, left, top, width, height, text, size=19, color=theme.text)
        return height + 0.12
    if kind in {"bullets", "numbered"}:
        items = list(block.get("items") or [])
        height = min(3.9, max(0.62, 0.46 * len(items)))
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        _add_bullets(shape.text_frame, items, numbered=kind == "numbered", font_size=19)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = theme.text
        return height + 0.12
    if kind == "quote":
        height = 1.0
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = theme.subtle
        box.line.fill.background()
        _add_textbox(slide, left + 0.25, top + 0.18, width - 0.5, 0.65, str(block.get("text") or ""), size=19, color=theme.text)
        return height + 0.16
    raise PresentationArtifactError(f"不支持的内容块：{kind}")


def _add_table_slide(
    slide, spec: dict[str, Any], number: int, theme: PresentationTheme
) -> None:
    _add_slide_header(slide, str(spec.get("title") or ""), number, theme)
    headers = [str(value) for value in spec.get("headers") or []]
    rows = [[str(value) for value in row] for row in spec.get("rows") or []]
    columns = len(headers) or (len(rows[0]) if rows else 0)
    if not columns or any(len(row) != columns for row in rows):
        raise PresentationArtifactError("表格列数不一致")
    table_shape = slide.shapes.add_table(
        len(rows) + (1 if headers else 0), columns,
        Inches(0.75), Inches(1.55), Inches(11.85), Inches(4.95),
    )
    table = table_shape.table
    values = ([headers] if headers else []) + rows
    for row_index, values_row in enumerate(values):
        for column_index, value in enumerate(values_row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.primary if headers and row_index == 0 else (theme.subtle if row_index % 2 == 0 else theme.surface)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_run_font(run, 16, bold=headers and row_index == 0, color=theme.on_primary if headers and row_index == 0 else theme.text)


def _first_value(spec: dict[str, Any], *keys: str) -> Any:
    for key in keys:
        value = spec.get(key)
        if value not in (None, "", []):
            return value
    return None


def _normalize_blocks(value: Any) -> list[dict[str, Any]]:
    if value in (None, "", []):
        return []
    if isinstance(value, str):
        return [{"type": "paragraph", "text": value}]
    if isinstance(value, dict):
        if value.get("type"):
            return [dict(value)]
        nested = _first_value(value, "blocks", "content", "body", "text")
        return _normalize_blocks(nested)
    if isinstance(value, list):
        if all(isinstance(item, dict) and item.get("type") for item in value):
            return [dict(item) for item in value]
        return [{"type": "bullets", "items": value}]
    return [{"type": "paragraph", "text": str(value)}]


def _normalize_slide_spec(raw: dict[str, Any]) -> dict[str, Any]:
    spec = dict(raw)
    explicit_kind = any(spec.get(key) not in (None, "") for key in ("type", "layout", "kind"))
    requested_kind = str(_first_value(spec, "type", "layout", "kind") or "content").lower()
    kind_aliases = {
        "title": "section",
        "title_slide": "section",
        "title-only": "section",
        "title_only": "section",
        "title_and_content": "content",
        "title-and-content": "content",
        "bullet": "content",
        "bullets": "content",
        "text": "content",
        "two-column": "two_column",
        "two columns": "two_column",
        "comparison": "two_column",
        "big_number": "metric",
        "number": "metric",
        "kpi": "metric",
        "steps": "process",
        "timeline": "process",
        "takeaway": "statement",
        "headline": "statement",
    }
    kind = kind_aliases.get(requested_kind, requested_kind)
    if "headers" in spec or "rows" in spec:
        kind = "table"
    elif "left" in spec or "right" in spec:
        kind = "two_column"
    elif any(key in spec for key in ("value", "metric")):
        kind = "metric"
    elif "steps" in spec:
        kind = "process"
    elif "statement" in spec:
        kind = "statement"
    candidate_title = str(_first_value(spec, "title", "heading", "name") or "")
    if not explicit_kind and kind == "content":
        title_lower = candidate_title.lower()
        raw_body = _first_value(spec, "content", "body", "text")
        raw_items = _first_value(spec, "bullets", "points", "items")
        if raw_items not in (None, "", []) and any(
            keyword in title_lower
            for keyword in ("流程", "步骤", "路径", "阶段", "路线图", "计划", "roadmap", "process")
        ):
            kind = "process"
            spec["steps"] = raw_items
        elif isinstance(raw_body, str) and len(raw_body.strip()) <= 220 and any(
            keyword in title_lower
            for keyword in ("结论", "核心", "观点", "建议", "判断", "启示", "takeaway", "recommendation")
        ):
            kind = "statement"
            spec["statement"] = raw_body
        elif isinstance(raw_body, str) and any(
            keyword in title_lower
            for keyword in ("增长", "占比", "转化", "收入", "成本", "利润", "规模", "指标", "率", "kpi")
        ):
            match = re.search(r"(?:[¥￥$])?\d+(?:\.\d+)?(?:%|万|亿|倍|[xX])?", raw_body)
            if match:
                kind = "metric"
                spec["_inferred_value"] = match.group(0)
                spec["_inferred_context"] = raw_body
    spec["type"] = kind
    spec["title"] = candidate_title
    spec["subtitle"] = str(_first_value(spec, "subtitle", "description") or "")

    if kind == "content":
        blocks = _normalize_blocks(spec.get("blocks"))
        if not blocks:
            blocks = _normalize_blocks(_first_value(spec, "content", "body", "text"))
        bullets = _first_value(spec, "bullets", "points", "items")
        if not blocks and bullets not in (None, "", []):
            blocks.extend([{"type": "bullets", "items": list(bullets) if isinstance(bullets, list) else [bullets]}])
        if not blocks and spec["subtitle"]:
            blocks = [{"type": "paragraph", "text": spec["subtitle"]}]
        spec["blocks"] = blocks
        if not spec["title"] and not blocks:
            raise PresentationArtifactError("正文幻灯片缺少标题和内容")
    elif kind == "statement":
        spec["text"] = str(_first_value(spec, "text", "content", "body", "statement") or "")
        if not spec["title"] and not spec["text"]:
            raise PresentationArtifactError("观点幻灯片缺少内容")
    elif kind == "metric":
        spec["value"] = str(_first_value(spec, "value", "metric", "number", "_inferred_value") or "")
        spec["label"] = str(_first_value(spec, "label", "caption", "subtitle") or "")
        spec["context"] = str(_first_value(spec, "context", "description", "body", "text", "_inferred_context") or "")
        if not spec["value"]:
            raise PresentationArtifactError("大数字幻灯片缺少 value")
    elif kind == "process":
        steps = _first_value(spec, "steps", "items", "points", "bullets") or []
        spec["steps"] = steps if isinstance(steps, list) else [steps]
        if not spec["steps"]:
            raise PresentationArtifactError("流程幻灯片缺少 steps")
    elif kind == "section":
        if not spec["title"] and not spec["subtitle"]:
            raise PresentationArtifactError("分节幻灯片缺少标题")
    elif kind == "two_column":
        for side in ("left", "right"):
            column = dict(spec.get(side) or {})
            column["heading"] = str(_first_value(column, "heading", "title", "name") or "")
            column_items = _first_value(column, "items", "bullets", "points")
            if column_items is not None:
                column["items"] = column_items if isinstance(column_items, list) else [column_items]
            else:
                column["text"] = str(_first_value(column, "text", "content", "body") or "")
            spec[side] = column
        if not spec["title"] and not any(
            spec[side].get("heading") or spec[side].get("text") or spec[side].get("items")
            for side in ("left", "right")
        ):
            raise PresentationArtifactError("双栏幻灯片缺少内容")
    elif kind != "table":
        raise PresentationArtifactError(f"不支持的幻灯片类型：{requested_kind}")
    return spec


def _add_highlight_list(
    slide, items: list[Any], *, numbered: bool, theme: PresentationTheme
) -> None:
    count = len(items)
    row_height = min(0.93, 4.75 / max(count, 1))
    top = 1.58
    for index, item in enumerate(items):
        text = str(item.get("text") or "") if isinstance(item, dict) else str(item)
        marker = str(index + 1).zfill(2) if numbered else "•"
        _add_textbox(
            slide, 0.82, top, 0.72, row_height, marker,
            size=24, bold=True, color=theme.accent, vertical=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide, 1.72, top, 10.45, row_height, text,
            size=22 if count <= 4 else 19, color=theme.text,
            vertical=MSO_ANCHOR.MIDDLE,
        )
        if index < count - 1:
            separator = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1.72), Inches(top + row_height),
                Inches(10.45), Inches(0.018),
            )
            separator.fill.solid()
            separator.fill.fore_color.rgb = theme.subtle
            separator.line.fill.background()
        top += row_height + 0.11


def _add_statement_slide(
    slide, spec: dict[str, Any], number: int, theme: PresentationTheme
) -> None:
    title = str(spec.get("title") or "核心观点")
    text = str(spec.get("text") or "")
    _add_textbox(slide, 0.82, 0.58, 10.7, 0.45, title.upper(), size=16, bold=True, color=theme.accent)
    _add_textbox(
        slide, 0.82, 1.55, 10.95, 3.25, text or title,
        size=38 if len(text or title) <= 55 else 30, bold=True, color=theme.text,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(5.42), Inches(2.15), Inches(0.09)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.accent
    line.line.fill.background()
    _add_textbox(slide, 12.05, 6.65, 0.62, 0.3, f"{number:02d}", size=11, color=theme.muted, align=PP_ALIGN.RIGHT)


def _add_metric_slide(
    slide, spec: dict[str, Any], number: int, theme: PresentationTheme
) -> None:
    _add_slide_header(slide, str(spec.get("title") or "关键指标"), number, theme)
    value = str(spec.get("value") or "")
    _add_textbox(slide, 0.82, 1.72, 7.0, 2.0, value, size=64, bold=True, color=theme.accent, vertical=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 0.9, 3.72, 6.9, 0.58, str(spec.get("label") or ""), size=24, bold=True, color=theme.text)
    context = str(spec.get("context") or "")
    if context:
        _add_textbox(slide, 7.55, 2.15, 4.55, 2.25, context, size=20, color=theme.text, vertical=MSO_ANCHOR.MIDDLE)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.15), Inches(2.08), Inches(0.08), Inches(2.42))
        accent.fill.solid()
        accent.fill.fore_color.rgb = theme.subtle
        accent.line.fill.background()


def _add_process_slide(
    slide, spec: dict[str, Any], number: int, theme: PresentationTheme
) -> None:
    _add_slide_header(slide, str(spec.get("title") or "流程"), number, theme)
    steps = list(spec.get("steps") or [])[:6]
    count = len(steps)
    start_x, end_x, y = 1.25, 12.0, 3.05
    if count > 1:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(y + 0.28),
            Inches(end_x - start_x), Inches(0.055),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme.subtle
        line.line.fill.background()
    gap = (end_x - start_x) / max(count - 1, 1)
    for index, step in enumerate(steps):
        if isinstance(step, dict):
            heading = str(_first_value(step, "title", "heading", "text") or "")
            detail = str(_first_value(step, "description", "detail", "body") or "")
        else:
            heading, detail = str(step), ""
        x = start_x + gap * index
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.28), Inches(y), Inches(0.58), Inches(0.58))
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme.accent
        circle.line.fill.background()
        _add_textbox(slide, x - 0.25, y + 0.08, 0.5, 0.3, str(index + 1), size=14, bold=True, color=theme.on_primary, align=PP_ALIGN.CENTER)
        width = min(2.15, gap * 0.86) if count > 1 else 4.0
        _add_textbox(slide, x - width / 2, y + 0.82, width, 0.55, heading, size=18, bold=True, color=theme.text, align=PP_ALIGN.CENTER)
        if detail:
            _add_textbox(slide, x - width / 2, y + 1.42, width, 0.86, detail, size=15, color=theme.muted, align=PP_ALIGN.CENTER)


def _continued_title(title: str, index: int) -> str:
    return title if index == 0 else f"{title}（续）"


def _expand_slide_specs(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    expanded: list[dict[str, Any]] = []
    for raw in slides:
        spec = _normalize_slide_spec(dict(raw))
        kind = spec["type"]
        if kind == "content":
            blocks = list(spec.get("blocks") or [])
            if len(blocks) == 1 and blocks[0].get("type") in {"bullets", "numbered"}:
                items = list(blocks[0].get("items") or [])
                if len(items) > 5:
                    for index, start in enumerate(range(0, len(items), 5)):
                        part = dict(spec)
                        part["title"] = _continued_title(spec["title"], index)
                        part["blocks"] = [{**blocks[0], "items": items[start:start + 5]}]
                        expanded.append(part)
                    continue
        if kind == "table":
            rows = list(spec.get("rows") or [])
            if len(rows) > 7:
                for index, start in enumerate(range(0, len(rows), 7)):
                    part = dict(spec)
                    part["title"] = _continued_title(spec["title"], index)
                    part["rows"] = rows[start:start + 7]
                    expanded.append(part)
                continue
        if kind == "process":
            steps = list(spec.get("steps") or [])
            if len(steps) > 5:
                for index, start in enumerate(range(0, len(steps), 5)):
                    part = dict(spec)
                    part["title"] = _continued_title(spec["title"], index)
                    part["steps"] = steps[start:start + 5]
                    expanded.append(part)
                continue
        expanded.append(spec)
    return expanded


def _validate_created_presentation(output_path: Path, expected_slides: int) -> None:
    try:
        presentation = Presentation(output_path)
    except Exception as error:
        raise PresentationArtifactError("生成的 PPTX 无法重新打开") from error
    if len(presentation.slides) != expected_slides:
        raise PresentationArtifactError("生成的 PPTX 页数校验失败")
    for index, slide in enumerate(presentation.slides, start=1):
        visible_text = [
            shape.text.strip() for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if not visible_text:
            raise PresentationArtifactError(f"第 {index} 页没有可见文字")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise PresentationArtifactError(f"第 {index} 页存在越界元素")
            if shape.left + shape.width > presentation.slide_width + Inches(0.02):
                raise PresentationArtifactError(f"第 {index} 页存在横向越界元素")
            if shape.top + shape.height > presentation.slide_height + Inches(0.02):
                raise PresentationArtifactError(f"第 {index} 页存在纵向越界元素")


def _append_slide(
    presentation: Presentation,
    spec: dict[str, Any],
    number: int,
    theme: PresentationTheme | None = None,
) -> None:
    theme = theme or THEMES["business"]
    spec = _normalize_slide_spec(spec)
    kind = str(spec.get("type") or "content")
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = theme.background

    if kind == "section":
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        band.fill.solid()
        band.fill.fore_color.rgb = theme.primary
        band.line.fill.background()
        _add_textbox(slide, 1.0, 2.45, 11.3, 1.0, str(spec.get("title") or ""), size=42, bold=True, color=theme.on_primary, align=PP_ALIGN.CENTER)
        _add_textbox(slide, 1.3, 3.62, 10.7, 0.62, str(spec.get("subtitle") or ""), size=20, color=theme.subtle, align=PP_ALIGN.CENTER)
        return
    if kind == "table":
        _add_table_slide(slide, spec, number, theme)
        return
    if kind == "statement":
        _add_statement_slide(slide, spec, number, theme)
        return
    if kind == "metric":
        _add_metric_slide(slide, spec, number, theme)
        return
    if kind == "process":
        _add_process_slide(slide, spec, number, theme)
        return

    _add_slide_header(slide, str(spec.get("title") or ""), number, theme)
    if kind == "two_column":
        columns = [dict(spec.get("left") or {}), dict(spec.get("right") or {})]
        for index, column in enumerate(columns):
            left = 0.78 + index * 6.05
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.55), Inches(5.65), Inches(4.95))
            panel.fill.solid()
            panel.fill.fore_color.rgb = theme.surface
            panel.line.fill.background()
            panel.line.color.rgb = theme.subtle
            _add_textbox(slide, left + 0.3, 1.84, 5.05, 0.52, str(column.get("heading") or ""), size=24, bold=True, color=theme.accent)
            body = slide.shapes.add_textbox(Inches(left + 0.3), Inches(2.48), Inches(5.05), Inches(3.62))
            items = column.get("items")
            if items is not None:
                _add_bullets(body.text_frame, list(items), font_size=18)
                for paragraph in body.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = theme.text
            else:
                body.text_frame.text = str(column.get("text") or "")
                for paragraph in body.text_frame.paragraphs:
                    for run in paragraph.runs:
                        _set_run_font(run, 19, color=theme.text)
        return
    blocks = list(spec.get("blocks") or [])
    if len(blocks) == 1 and str(blocks[0].get("type") or "") in {"bullets", "numbered"}:
        items = list(blocks[0].get("items") or [])
        if 2 <= len(items) <= 5:
            _add_highlight_list(
                slide, items,
                numbered=str(blocks[0].get("type")) == "numbered",
                theme=theme,
            )
            return
    top = 1.55
    for block in blocks:
        top += _add_content_block(slide, dict(block), 0.8, top, 11.75, theme)
        if top > 6.75:
            raise PresentationArtifactError(f"幻灯片“{spec.get('title') or ''}”内容过多")


def create_presentation(
    output_path: Path,
    title: str,
    subtitle: str,
    slides: list[dict[str, Any]],
    theme_name: str = "auto",
    design_profile: dict[str, Any] | None = None,
) -> str:
    if not title.strip():
        raise PresentationArtifactError("PPT 标题不能为空")
    expanded_slides = _expand_slide_specs(slides)
    if len(expanded_slides) > 49:
        raise PresentationArtifactError("PPT 自动拆页后最多支持 50 页")
    theme = (
        select_theme_for_content(title, subtitle, slides, design_profile)
        if str(theme_name or "auto").strip().lower() == "auto"
        else resolve_theme(theme_name)
    )
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    cover = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = cover.background.fill
    background.solid()
    background.fore_color.rgb = theme.primary
    accent = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.25), Inches(0.16), Inches(4.35))
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme.accent
    accent.line.fill.background()
    title_size = 50 if len(title) <= 30 else 44
    _add_textbox(cover, 1.18, 1.9, 10.8, 1.7, title, size=title_size, bold=True, color=theme.on_primary, vertical=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_textbox(cover, 1.2, 3.78, 10.6, 0.72, subtitle, size=22, color=theme.subtle)
    for number, spec in enumerate(expanded_slides, start=2):
        _append_slide(presentation, dict(spec), number, theme)
    presentation.core_properties.title = title
    presentation.core_properties.author = "Nano Agent"
    presentation.save(output_path)
    _validate_created_presentation(output_path, len(expanded_slides) + 1)
    return theme.name


def _delete_slide(presentation: Presentation, one_based_number: int) -> None:
    if one_based_number < 1 or one_based_number > len(presentation.slides):
        raise PresentationArtifactError(f"幻灯片页码不存在：{one_based_number}")
    slide_id = presentation.slides._sldIdLst[one_based_number - 1]
    relationship_id = slide_id.rId
    presentation.part.drop_rel(relationship_id)
    presentation.slides._sldIdLst.remove(slide_id)


def _replace_in_text_frame(frame, old: str, new: str) -> int:
    count = 0
    for paragraph in frame.paragraphs:
        if old not in paragraph.text:
            continue
        if any(old in run.text for run in paragraph.runs):
            for run in paragraph.runs:
                occurrences = run.text.count(old)
                run.text = run.text.replace(old, new)
                count += occurrences
        else:
            occurrences = paragraph.text.count(old)
            paragraph.text = paragraph.text.replace(old, new)
            count += occurrences
    return count


def edit_presentation(
    source_path: Path,
    output_path: Path,
    operations: list[dict[str, Any]],
) -> dict[str, int]:
    presentation = Presentation(source_path)
    counts = {"replacements": 0, "deleted_slides": 0, "appended_slides": 0}
    for operation in operations:
        kind = str(operation.get("type") or "")
        if kind == "replace_text":
            old = str(operation.get("old") or "")
            if not old:
                raise PresentationArtifactError("replace_text.old 不能为空")
            new = str(operation.get("new") or "")
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        counts["replacements"] += _replace_in_text_frame(
                            shape.text_frame, old, new
                        )
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                counts["replacements"] += _replace_in_text_frame(
                                    cell.text_frame, old, new
                                )
        elif kind == "delete_slide":
            _delete_slide(presentation, int(operation.get("slide_number") or 0))
            counts["deleted_slides"] += 1
        elif kind == "append_slides":
            new_slides = _expand_slide_specs(list(operation.get("slides") or []))
            for spec in new_slides:
                _append_slide(presentation, dict(spec), len(presentation.slides) + 1)
            counts["appended_slides"] += len(new_slides)
        else:
            raise PresentationArtifactError(f"不支持的编辑操作：{kind}")
    presentation.save(output_path)
    return counts
