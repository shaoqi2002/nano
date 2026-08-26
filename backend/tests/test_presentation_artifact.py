import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from app.service.presentation_artifact import (
    _normalize_slide_spec,
    create_presentation,
    edit_presentation,
    presentation_text,
    select_theme_for_content,
)


class PresentationArtifactTests(unittest.TestCase):
    def test_creates_structured_widescreen_presentation(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "generated.pptx"
            create_presentation(
                output,
                "Nano 演示",
                "聊天中直接生成",
                [
                    {
                        "type": "content",
                        "title": "核心能力",
                        "blocks": [
                            {"type": "heading", "text": "聊天式工作流"},
                            {"type": "bullets", "items": ["创建 PPT", "编辑附件"]},
                        ],
                    },
                    {
                        "type": "table",
                        "title": "功能清单",
                        "headers": ["功能", "状态"],
                        "rows": [["生成", "可用"], ["编辑", "可用"]],
                    },
                ],
            )
            deck = Presentation(output)

            self.assertEqual(len(deck.slides), 3)
            self.assertGreater(deck.slide_width, deck.slide_height)
            extracted = presentation_text(output.read_bytes())
            self.assertIn("Nano 演示", extracted)
            self.assertIn("编辑附件", extracted)
            self.assertIn("功能\t状态", extracted)

    def test_edits_copy_without_modifying_source(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pptx"
            output = root / "edited.pptx"
            create_presentation(
                source,
                "原始标题",
                "",
                [{"type": "content", "title": "旧内容", "blocks": []}],
            )
            summary = edit_presentation(
                source,
                output,
                [
                    {"type": "replace_text", "old": "旧内容", "new": "新内容"},
                    {
                        "type": "append_slides",
                        "slides": [{"type": "section", "title": "附录"}],
                    },
                ],
            )

            self.assertEqual(summary["replacements"], 1)
            self.assertEqual(summary["appended_slides"], 1)
            self.assertIn("旧内容", presentation_text(source.read_bytes()))
            self.assertIn("新内容", presentation_text(output.read_bytes()))
            self.assertEqual(len(Presentation(output).slides), 3)

    def test_accepts_simplified_agent_slide_fields(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "simple.pptx"
            create_presentation(
                output,
                "简化输入",
                "",
                [
                    {"heading": "摘要", "content": "这段内容不使用 blocks 字段。"},
                    {"title": "重点", "bullets": ["第一点", "第二点"]},
                    {
                        "layout": "title_and_content",
                        "title": "结论",
                        "body": ["结论一", "结论二"],
                    },
                ],
            )

            extracted = presentation_text(output.read_bytes())
            self.assertEqual(len(Presentation(output).slides), 4)
            self.assertIn("这段内容不使用 blocks 字段", extracted)
            self.assertIn("第一点", extracted)
            self.assertIn("结论一", extracted)

    def test_rejects_completely_empty_slide(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            with self.assertRaisesRegex(Exception, "缺少标题和内容"):
                create_presentation(
                    Path(directory) / "empty.pptx", "标题", "", [{}]
                )

    def test_supports_themes_visual_layouts_and_automatic_splitting(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "visual.pptx"
            create_presentation(
                output,
                "视觉升级",
                "modern 主题",
                [
                    {"type": "statement", "title": "核心判断", "text": "每一页只表达一个关键观点。"},
                    {"type": "metric", "title": "增长", "value": "42%", "label": "年度增长率", "context": "增长主要来自企业客户。"},
                    {"type": "process", "title": "实施路径", "steps": ["规划", "设计", "生成", "检查"]},
                    {"title": "九项能力", "bullets": [f"能力 {index}" for index in range(1, 10)]},
                ],
                theme_name="modern",
            )

            deck = Presentation(output)
            extracted = presentation_text(output.read_bytes())
            self.assertEqual(len(deck.slides), 6)
            self.assertIn("42%", extracted)
            self.assertIn("九项能力（续）", extracted)

    def test_automatically_selects_theme_from_deck_content(self) -> None:
        first = select_theme_for_content("AI 智能体技术架构", "", [])
        repeated = select_theme_for_content("AI 智能体技术架构", "", [])
        different = select_theme_for_content("品牌创意与用户体验", "", [])
        branded = select_theme_for_content(
            "品牌发布",
            "",
            [],
            {
                "mood": "克制而有力量",
                "mode": "dark",
                "primary_color": "#112233",
                "accent_color": "#FF8800",
            },
        )

        self.assertTrue(first.name.startswith("adaptive-"))
        self.assertEqual(first, repeated)
        self.assertNotEqual(first.accent, different.accent)
        self.assertEqual(tuple(branded.primary), (0x11, 0x22, 0x33))
        self.assertEqual(tuple(branded.accent), (0xFF, 0x88, 0x00))
        with self.assertRaisesRegex(Exception, "#RRGGBB"):
            select_theme_for_content(
                "无效品牌色", "", [], {"primary_color": "blue"}
            )

    def test_automatically_infers_layout_from_slide_content(self) -> None:
        statement = _normalize_slide_spec({
            "title": "核心结论",
            "content": "增长来自高价值企业客户。",
        })
        metric = _normalize_slide_spec({
            "title": "年度增长率",
            "content": "用户规模同比增长 42%，主要来自企业客户。",
        })
        process = _normalize_slide_spec({
            "title": "实施步骤",
            "bullets": ["规划", "设计", "上线"],
        })
        explicit = _normalize_slide_spec({
            "type": "content",
            "title": "核心结论",
            "content": "保持普通正文布局。",
        })

        self.assertEqual(statement["type"], "statement")
        self.assertEqual(metric["type"], "metric")
        self.assertEqual(metric["value"], "42%")
        self.assertEqual(process["type"], "process")
        self.assertEqual(explicit["type"], "content")


if __name__ == "__main__":
    unittest.main()
