import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
from uuid import UUID

from pptx import Presentation

from app.service.chat_artifact import (
    ChatArtifactNotFoundError,
    get_chat_artifact,
    store_chat_artifact,
    store_chat_artifact_bytes,
)
from app.tools.local_write import (
    presentation_create,
    presentation_edit_attachment,
    word_create_document,
    word_edit_document,
)


class ChatArtifactTests(unittest.TestCase):
    def test_stores_and_resolves_download_without_document_library(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.docx"
            source.write_bytes(b"docx-content")
            artifact_root = root / "artifacts"
            with patch("app.service.chat_artifact.CHAT_ARTIFACT_DIR", artifact_root):
                result = store_chat_artifact(source, "report.docx")
                stored, media_type = get_chat_artifact(result["artifact_id"])

            self.assertEqual(stored.read_bytes(), b"docx-content")
            self.assertEqual(stored.name, "report.docx")
            self.assertIn("wordprocessingml.document", media_type)
            self.assertEqual(
                result["download_url"],
                f"/api/artifacts/{result['artifact_id']}?workspace_id=00000000-0000-0000-0000-0000000000c4",
            )

    def test_artifacts_are_isolated_by_workspace(self) -> None:
        first = UUID("00000000-0000-0000-0000-000000000111")
        second = UUID("00000000-0000-0000-0000-000000000222")
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.docx"
            source.write_bytes(b"private")
            with patch("app.service.chat_artifact.CHAT_ARTIFACT_DIR", root / "artifacts"):
                result = store_chat_artifact(source, "private.docx", first)
                with self.assertRaises(ChatArtifactNotFoundError):
                    get_chat_artifact(result["artifact_id"], second)
                stored, _ = get_chat_artifact(result["artifact_id"], first)
                self.assertEqual(stored.read_bytes(), b"private")

    def test_missing_artifact_is_rejected(self) -> None:
        with tempfile.TemporaryDirectory() as directory, patch(
            "app.service.chat_artifact.CHAT_ARTIFACT_DIR", Path(directory)
        ):
            with self.assertRaises(ChatArtifactNotFoundError):
                get_chat_artifact("00000000-0000-0000-0000-000000000000")

    def test_word_tools_default_to_docx_only(self) -> None:
        self.assertFalse(word_create_document.args_schema.model_fields["create_pdf"].default)
        self.assertFalse(word_edit_document.args_schema.model_fields["create_pdf"].default)


class ChatArtifactToolTests(unittest.IsolatedAsyncioTestCase):
    async def test_word_create_returns_one_direct_download(self) -> None:
        with tempfile.TemporaryDirectory() as directory, patch(
            "app.service.chat_artifact.CHAT_ARTIFACT_DIR", Path(directory)
        ):
            result = await word_create_document.ainvoke({
                "filename": "report.docx",
                "title": "Report",
                "blocks": [{"type": "paragraph", "text": "Direct download"}],
            })

        self.assertEqual(len(result["outputs"]), 1)
        output = result["outputs"][0]
        self.assertEqual(output["kind"], "word")
        self.assertEqual(output["filename"], "report.docx")
        self.assertTrue(output["download_url"].startswith("/api/artifacts/"))
        self.assertNotIn("document_id", output)

    async def test_presentation_create_and_chat_attachment_edit_return_downloads(self) -> None:
        with tempfile.TemporaryDirectory() as directory, patch(
            "app.service.chat_artifact.CHAT_ARTIFACT_DIR", Path(directory)
        ):
            created = await presentation_create.ainvoke({
                "filename": "briefing.pptx",
                "title": "旧标题",
                "slides": [{
                    "type": "content",
                    "title": "计划",
                    "blocks": [{"type": "paragraph", "text": "第一版"}],
                }],
            })
            created_output = created["outputs"][0]
            source_path, _ = get_chat_artifact(created_output["artifact_id"])
            source = store_chat_artifact_bytes(source_path.read_bytes(), "uploaded.pptx")
            edited = await presentation_edit_attachment.ainvoke({
                "artifact_id": source["artifact_id"],
                "output_filename": "briefing-v2.pptx",
                "operations": [{"type": "replace_text", "old": "旧标题", "new": "新标题"}],
            })
            edited_path, _ = get_chat_artifact(edited["outputs"][0]["artifact_id"])
            edited_deck = Presentation(edited_path)
            edited_text = "\n".join(
                shape.text for slide in edited_deck.slides
                for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            )

        self.assertEqual(created_output["kind"], "presentation")
        self.assertTrue(created["selected_theme"].startswith("adaptive-"))
        self.assertEqual(edited["outputs"][0]["kind"], "presentation")
        self.assertEqual(edited["edit_summary"]["replacements"], 1)
        self.assertEqual(edited_deck.core_properties.title, "旧标题")
        self.assertIn("新标题", edited_text)


if __name__ == "__main__":
    unittest.main()
